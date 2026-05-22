from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
REPORTS = ROOT / "reports"
FIGURES = REPORTS / "figures"
OUTPUT = REPORTS / "Фамилия_Имя_ML_прогнозирование_спроса_инструменты.pptx"

BG = RGBColor(248, 250, 252)
INK = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(13, 148, 136)
RED = RGBColor(220, 38, 38)


def add_title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = INK
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.88), Inches(12), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(11)
        sp.font.color.rgb = MUTED


def add_footer(slide, number: int):
    box = slide.shapes.add_textbox(Inches(11.9), Inches(7.05), Inches(0.8), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, items, left, top, width, height, size=15):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(8)


def add_metric_card(slide, label, value, left, top, color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(2.25), Inches(0.95))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(226, 232, 240)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(9)
    p2.font.color.rgb = MUTED
    p2.alignment = PP_ALIGN.CENTER


def set_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def main():
    metrics = json.loads((REPORTS / "metrics.json").read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Автоматизированная ML-система прогнозирования спроса", "Ниша: ручные строительные инструменты, Building Supply Store Sales 2020-2024")
    add_bullets(
        slide,
        [
            "Цель: прогноз недельного спроса по товару и магазину",
            "Бизнес-эффект: точнее закупки, меньше дефицита и лишних остатков",
            "Результат: воспроизводимый ML-пайплайн с AutoML, Docker, CI/CD, тестами и мониторингом",
        ],
        0.75,
        1.55,
        5.8,
        2.2,
        17,
    )
    slide.shapes.add_picture(str(FIGURES / "weekly_tool_demand.png"), Inches(6.45), Inches(1.35), Inches(6.25), Inches(3.15))
    add_metric_card(slide, "рядов товар-магазин", "77", 0.75, 4.75, GREEN)
    add_metric_card(slide, "строк после ETL", "20 174", 3.2, 4.75, BLUE)
    add_metric_card(slide, "holdout", "16 недель", 5.65, 4.75, RED)
    add_footer(slide, 1)

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Данные и ETL", "Из транзакций продаж построены недельные временные ряды")
    add_bullets(
        slide,
        [
            "Extract: чтение CSV и проверка обязательных колонок",
            "Transform: фильтрация product_type = Tool, недельная агрегация, полный календарь недель",
            "Features: календарные признаки, лаги 1/2/4/8 недель, rolling mean/std, цена и delivery share",
            "Load: joblib-модель, CSV-прогнозы, JSON-метрики, PNG-графики",
        ],
        0.7,
        1.35,
        5.9,
        3.6,
        15,
    )
    slide.shapes.add_picture(str(FIGURES / "top_tool_products.png"), Inches(6.65), Inches(1.25), Inches(5.95), Inches(3.25))
    add_footer(slide, 2)

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Архитектура ML-системы", "Автоматизация выбора модели через GridSearchCV и TimeSeriesSplit")
    add_bullets(
        slide,
        [
            "Preprocessing: OneHotEncoder для категорий, StandardScaler для числовых признаков",
            "AutoML-кандидаты: Ridge, RandomForestRegressor, HistGradientBoostingRegressor",
            "Валидация: временные разбиения без утечки будущих данных",
            f"Лучшая модель: {metrics['best_estimator']}",
        ],
        0.85,
        1.4,
        11.4,
        3.3,
        16,
    )
    add_metric_card(slide, "max_depth", "8", 1.0, 5.05, BLUE)
    add_metric_card(slide, "min_samples_leaf", "3", 3.55, 5.05, GREEN)
    add_metric_card(slide, "n_estimators", "160", 6.1, 5.05, RED)
    add_metric_card(slide, "CV MAE", "2.603", 8.65, 5.05, BLUE)
    add_footer(slide, 3)

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Результаты модели", "Holdout: последние 16 недель")
    m = metrics["metrics"]
    add_metric_card(slide, "MAE", f"{m['mae']:.3f}", 0.8, 1.35, BLUE)
    add_metric_card(slide, "RMSE", f"{m['rmse']:.3f}", 3.25, 1.35, RED)
    add_metric_card(slide, "SMAPE", f"{m['smape']:.1f}%", 5.7, 1.35, GREEN)
    add_metric_card(slide, "R2", f"{m['r2']:.3f}", 8.15, 1.35, BLUE)
    slide.shapes.add_picture(str(FIGURES / "forecast_vs_actual.png"), Inches(0.85), Inches(2.65), Inches(5.9), Inches(3.15))
    slide.shapes.add_picture(str(FIGURES / "residuals.png"), Inches(6.95), Inches(2.65), Inches(5.45), Inches(3.15))
    add_footer(slide, 4)

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Тестирование, Docker и CI/CD", "Проект готов к воспроизводимому запуску")
    add_bullets(
        slide,
        [
            "Тесты pytest: ETL, признаки, набор features, SMAPE и regression metrics",
            "Локальный результат: 7 passed",
            "Docker: python:3.11-slim, непривилегированный пользователь appuser, запуск train-команды",
            "GitHub Actions: установка зависимостей, pytest, сборка Docker image",
        ],
        0.85,
        1.35,
        10.9,
        4.0,
        16,
    )
    add_footer(slide, 5)

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Мониторинг и выводы для бизнеса", "Контроль качества данных, дрейфа и времени обучения")
    add_bullets(
        slide,
        [
            "Data quality: пропуски 0% по подготовленным признакам",
            "Drift: PSI выявляет календарный сдвиг на holdout-периоде конца 2024 года",
            f"Training time: {metrics['training_seconds']} секунд",
            "Рекомендация: использовать прогноз для еженедельного пополнения запасов и переобучать модель ежемесячно",
            "Для редкого спроса дополнить модель правилом минимального страхового запаса",
        ],
        0.85,
        1.35,
        11.2,
        4.2,
        16,
    )
    add_footer(slide, 6)

    prs.save(OUTPUT)
    print(str(OUTPUT).encode("unicode_escape").decode("ascii"))


if __name__ == "__main__":
    main()
